from __future__ import annotations

import io
import os
from pathlib import Path

import numpy as np
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation as PptxPresentation
import pymupdf

from app.services.retrieval.chunker import chunk_text
from app.services.retrieval.rag_retriever import (
    EXPECTED_DIM,
    get_connection,
    get_model,
    vector_literal,
)


SUPPORTED_FILE_TYPES = {"pdf", "docx", "txt", "md", "xlsx", "pptx"}
MAX_FILE_SIZE = 20 * 1024 * 1024


def resolve_file_type(filename: str | None, file_type: str | None) -> str:
    if file_type:
        value = file_type.lower().strip().lstrip(".")
    elif filename:
        value = Path(filename).suffix.lower().lstrip(".")
    else:
        raise ValueError("file_type 또는 filename이 필요합니다.")

    if value not in SUPPORTED_FILE_TYPES:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {value}")

    return value


def extract_pdf_text(file_bytes: bytes) -> str:
    """PDF에서 텍스트를 추출한다."""
    try:
        doc = pymupdf.open(stream=file_bytes, filetype="pdf")
        pages = []

        for page in doc:
            page_text = page.get_text("text")
            if page_text:
                pages.append(page_text.strip())

        return "\n".join(pages).strip()

    except Exception as exc:
        raise ValueError(f"PDF 텍스트 추출에 실패했습니다: {exc}") from exc


def extract_docx_text(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_xlsx_text(file_bytes: bytes) -> str:
    """엑셀 각 시트를 순서대로 훑으며, 행마다 셀 값을 ' | '로 이어붙인 텍스트로 변환한다."""
    workbook = load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
    parts = []

    for sheet in workbook.worksheets:
        parts.append(f"[시트: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_pptx_text(file_bytes: bytes) -> str:
    prs = PptxPresentation(io.BytesIO(file_bytes))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_parts.append(shape.text_frame.text.strip())
        if slide_parts:
            parts.append(f"[슬라이드 {i}]\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def decode_text_file(file_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949"):
        try:
            return file_bytes.decode(encoding).strip()
        except UnicodeDecodeError:
            pass

    raise ValueError("TXT/MD 인코딩을 읽을 수 없습니다.")


def extract_text(file_bytes: bytes, file_type: str) -> str:
    if not file_bytes:
        raise ValueError("파일이 비어 있습니다.")

    if len(file_bytes) > MAX_FILE_SIZE:
        raise ValueError("파일 크기가 20MB를 초과했습니다.")

    if file_type == "pdf":
        text = extract_pdf_text(file_bytes)
    elif file_type == "docx":
        text = extract_docx_text(file_bytes)
    elif file_type == "xlsx":
        text = extract_xlsx_text(file_bytes)
    elif file_type == "pptx":
        text = extract_pptx_text(file_bytes)
    elif file_type in {"txt", "md"}:
        text = decode_text_file(file_bytes)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {file_type}")

    if not text.strip():
        raise ValueError("문서에서 텍스트를 추출하지 못했습니다.")

    return text.strip()


def embed_chunks(chunks: list[str]) -> np.ndarray:
    if not chunks:
        raise ValueError("embedding할 chunk가 없습니다.")

    model = get_model()

    try:
        batch_size = int(os.getenv("BGE_M3_BATCH_SIZE", "4"))
    except ValueError:
        batch_size = 4
    batch_size = max(1, min(batch_size, 64))

    output = model.encode(
        chunks,
        batch_size=batch_size,
        max_length=512,
        return_dense=True,
        return_sparse=False,
        return_colbert_vecs=False,
    )

    embeddings = np.asarray(
        output["dense_vecs"],
        dtype=np.float32,
    )

    if embeddings.shape != (len(chunks), EXPECTED_DIM):
        raise RuntimeError(
            f"잘못된 embedding shape: {embeddings.shape}"
        )

    if not np.isfinite(embeddings).all():
        raise RuntimeError("embedding에 NaN 또는 Inf가 있습니다.")

    norms = np.linalg.norm(embeddings, axis=1)

    if np.any(norms <= 0):
        raise RuntimeError("0 vector가 생성되었습니다.")

    return embeddings


def verify_document_owner(
    document_id: int,
    owner_user_id: int,
) -> None:
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT owner_user_id
                FROM documents
                WHERE id = :1
                """,
                (document_id,),
            )
            row = cur.fetchone()

    if row is None:
        raise ValueError(
            f"존재하지 않는 document_id입니다: {document_id}"
        )

    if row[0] != owner_user_id:
        raise PermissionError(
            "문서 소유자가 일치하지 않습니다."
        )


def load_document_metadata(
    document_id: int,
    owner_user_id: int,
) -> dict[str, str]:
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT title, document_type, description
                FROM documents
                WHERE id = :1
                  AND owner_user_id = :2
                """,
                (document_id, owner_user_id),
            )
            row = cur.fetchone()

    if row is None:
        raise ValueError(
            f"문서 메타데이터를 찾을 수 없습니다: {document_id}"
        )

    title, document_type, description = row
    return {
        "title": str(title or "").strip(),
        "document_type": str(document_type or "OTHER").strip(),
        "description": str(description or "").strip(),
    }


def build_embedding_inputs(
    chunks: list[str],
    metadata: dict[str, str],
) -> list[str]:
    """
    검색 embedding에는 파일명/유형/설명을 함께 넣되 DB의 실제 chunk content는
    원문 그대로 유지한다. 따라서 "전에 올린 이력서", "회사 보고서" 같은
    자연어 문서 찾기가 본문 단어에만 의존하지 않는다.
    """
    prefix_parts = []

    if metadata.get("title"):
        prefix_parts.append(f"제목: {metadata['title']}")
    if metadata.get("document_type"):
        prefix_parts.append(f"문서 유형: {metadata['document_type']}")
    if metadata.get("description"):
        prefix_parts.append(f"설명: {metadata['description']}")

    prefix = "\n".join(prefix_parts)

    if not prefix:
        return chunks

    return [
        f"{prefix}\n내용:\n{chunk}"
        for chunk in chunks
    ]


def save_chunk_texts(
    document_id: int,
    owner_user_id: int,
    chunks: list[str],
) -> None:
    """텍스트 추출 결과를 embedding과 독립적으로 먼저 보존한다."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT owner_user_id
                FROM documents
                WHERE id = :1
                FOR UPDATE
                """,
                (document_id,),
            )
            row = cur.fetchone()

            if row is None:
                raise ValueError(
                    f"존재하지 않는 document_id입니다: {document_id}"
                )

            if row[0] != owner_user_id:
                raise PermissionError(
                    "문서 소유자가 일치하지 않습니다."
                )

            cur.execute(
                "DELETE FROM document_chunks WHERE document_id = :1",
                (document_id,),
            )

            cur.executemany(
                """
                INSERT INTO document_chunks (
                    document_id, chunk_index, content, embedding
                )
                VALUES (:1, :2, :3, NULL)
                """,
                [
                    (document_id, chunk_index, content)
                    for chunk_index, content in enumerate(chunks)
                ],
            )

        conn.commit()


def save_chunk_embeddings(
    document_id: int,
    owner_user_id: int,
    embeddings: np.ndarray,
) -> None:
    """텍스트가 저장된 뒤 embedding만 별도로 채운다."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT owner_user_id
                FROM documents
                WHERE id = :1
                FOR UPDATE
                """,
                (document_id,),
            )
            row = cur.fetchone()

            if row is None:
                raise ValueError(
                    f"존재하지 않는 document_id입니다: {document_id}"
                )

            if row[0] != owner_user_id:
                raise PermissionError(
                    "문서 소유자가 일치하지 않습니다."
                )

            # PostgreSQL의 %s::vector 캐스트 → Oracle은 TO_VECTOR(:1)로 문자열을
            # VECTOR 타입으로 변환한다.
            cur.executemany(
                """
                UPDATE document_chunks
                SET embedding = TO_VECTOR(:1)
                WHERE document_id = :2
                  AND chunk_index = :3
                """,
                [
                    (
                        vector_literal(embeddings[chunk_index]),
                        document_id,
                        chunk_index,
                    )
                    for chunk_index in range(len(embeddings))
                ],
            )

        conn.commit()


def index_document(
    document_id: int,
    owner_user_id: int,
    file_bytes: bytes,
    filename: str | None = None,
    file_type: str | None = None,
) -> dict:
    resolved_file_type = resolve_file_type(
        filename=filename,
        file_type=file_type,
    )

    verify_document_owner(
        document_id=document_id,
        owner_user_id=owner_user_id,
    )

    text = extract_text(
        file_bytes=file_bytes,
        file_type=resolved_file_type,
    )

    chunks = chunk_text(
        text,
        min_chars=300,
        target_chars=400,
        max_chars=500,
    )

    if not chunks:
        raise ValueError("chunking 결과가 비어 있습니다.")

    # 가장 중요한 순서: 먼저 실제 텍스트를 저장한다. BGE-M3가 실패해도
    # 사용자는 방금 올린 문서의 전체 내용/요약을 읽을 수 있어야 한다.
    save_chunk_texts(
        document_id=document_id,
        owner_user_id=owner_user_id,
        chunks=chunks,
    )

    metadata = load_document_metadata(
        document_id=document_id,
        owner_user_id=owner_user_id,
    )

    embedding_error = None

    try:
        embeddings = embed_chunks(
            build_embedding_inputs(chunks, metadata)
        )
        save_chunk_embeddings(
            document_id=document_id,
            owner_user_id=owner_user_id,
            embeddings=embeddings,
        )
        status = "ready"
    except Exception as exc:
        # 텍스트는 이미 DB에 있으므로 "완전 실패"가 아니다.
        status = "text_ready"
        embedding_error = str(exc)[:1000]

    return {
        "document_id": document_id,
        "owner_user_id": owner_user_id,
        "file_type": resolved_file_type,
        "text_chars": len(text),
        "chunk_count": len(chunks),
        "embedding_dimension": EXPECTED_DIM,
        "status": status,
        "embedding_error": embedding_error,
    }
