from __future__ import annotations

import re
import textwrap
from io import BytesIO
from pathlib import Path

import fitz
from docx import Document as DocxDocument
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Pt


def _safe_filename(title: str) -> str:
    name = re.sub(r'[\\/:*?"<>|]+', "_", title).strip()
    return name or "document"


def _generate_docx(title: str, content: str) -> bytes:
    document = DocxDocument()
    document.add_heading(title, level=0)

    for line in content.splitlines():
        if line.strip():
            document.add_paragraph(line)
        else:
            document.add_paragraph("")

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _generate_pdf(title: str, content: str) -> bytes:
    pdf = fitz.open()

    page_width = 595
    page_height = 842
    margin = 50
    bottom_margin = 50

    page = pdf.new_page(width=page_width, height=page_height)
    y = 60

    def new_page():
        nonlocal page, y
        page = pdf.new_page(width=page_width, height=page_height)
        y = 60

    for title_line in textwrap.wrap(title, width=32) or [title]:
        page.insert_text(
            (margin, y),
            title_line,
            fontsize=17,
            fontname="korea",
        )
        y += 24

    y += 12

    for paragraph in content.splitlines():
        if not paragraph.strip():
            y += 10
            continue

        lines = textwrap.wrap(
            paragraph,
            width=48,
            break_long_words=True,
            replace_whitespace=False,
        ) or [""]

        for line in lines:
            if y >= page_height - bottom_margin:
                new_page()

            page.insert_text(
                (margin, y),
                line,
                fontsize=10.5,
                fontname="korea",
            )
            y += 17

        y += 4

    result = pdf.tobytes()
    pdf.close()
    return result



REPORT_FIELDS = [
    "문서명",
    "작성일",
    "작성부서",
    "작성자",
    "보고대상",
    "문서번호",
]

REPORT_SECTIONS = [
    "1. 보고 목적",
    "2. 주요 내용",
    "3. 현황 및 문제점",
    "4. 검토 내용",
    "5. 조치 / 실행 계획",
    "6. 요청 및 결정 사항",
    "첨부 / 참고 자료",
]


def _file_extension(filename: str | None) -> str:
    if not filename:
        return ""
    return Path(filename).suffix.lower().lstrip(".")


def _normalize_heading(text: str) -> str:
    text = re.sub(r"^\s*\d+\s*[.)]\s*", "", text.strip())
    return re.sub(r"\s+", " ", text).strip()


def _extract_template_text(
    template_bytes: bytes,
    template_filename: str | None,
) -> str:
    ext = _file_extension(template_filename)

    if ext == "pdf":
        pdf = fitz.open(
            stream=template_bytes,
            filetype="pdf",
        )
        try:
            return "\n".join(
                page.get_text("text")
                for page in pdf
            )
        finally:
            pdf.close()

    if ext == "docx":
        document = DocxDocument(BytesIO(template_bytes))
        values = []

        values.extend(
            p.text
            for p in document.paragraphs
            if p.text.strip()
        )

        for table in document.tables:
            for row in table.rows:
                values.append(
                    " | ".join(
                        cell.text.strip()
                        for cell in row.cells
                    )
                )

        return "\n".join(values)

    if ext == "xlsx":
        workbook = load_workbook(
            BytesIO(template_bytes),
            data_only=True,
        )

        values = []

        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                cells = [
                    str(cell.value).strip()
                    for cell in row
                    if cell.value not in (None, "")
                ]

                if cells:
                    values.append(" | ".join(cells))

        return "\n".join(values)

    if ext in {"md", "txt"}:
        return template_bytes.decode(
            "utf-8",
            errors="replace",
        )

    return ""


def _template_sections(template_text: str) -> list[str]:
    sections = []

    normalized_template = _normalize_heading(
        template_text
    )

    for section in REPORT_SECTIONS:
        if (
            _normalize_heading(section)
            in normalized_template
        ):
            sections.append(section)

    return sections or REPORT_SECTIONS


def _parse_report_content(
    content: str,
    sections: list[str],
) -> tuple[dict[str, str], dict[str, str]]:
    normalized_content = (
        content.replace("\r\n", "\n")
        .replace("\r", "\n")
    )

    for field in REPORT_FIELDS:
        normalized_content = re.sub(
            rf"(?<!\n)\s*(?={re.escape(field)}\s*[:：])",
            "\n",
            normalized_content,
        )

    for section in sorted(
        sections,
        key=len,
        reverse=True,
    ):
        normalized_content = re.sub(
            rf"(?<!\n)\s*(?={re.escape(section)}(?:\s|$))",
            "\n",
            normalized_content,
        )

    lines = [
        line.strip()
        for line in normalized_content.splitlines()
    ]

    fields = {}
    section_values = {}

    section_map = {
        _normalize_heading(section): section
        for section in sections
    }

    i = 0

    while i < len(lines):
        line = lines[i]

        if not line:
            i += 1
            continue

        field_found = False

        for field in REPORT_FIELDS:
            match = re.match(
                rf"^{re.escape(field)}\s*[:：]\s*(.*)$",
                line,
            )

            if match:
                fields[field] = match.group(1).strip()
                field_found = True
                break

        if field_found:
            i += 1
            continue

        normalized = _normalize_heading(line)
        section = section_map.get(normalized)

        if section:
            body = []
            j = i + 1

            while j < len(lines):
                candidate = lines[j]

                if (
                    _normalize_heading(candidate)
                    in section_map
                ):
                    break

                body.append(candidate)
                j += 1

            section_values[section] = (
                "\n".join(body).strip()
            )

            i = j
            continue

        i += 1

    return fields, section_values


def _is_report_template(template_text: str) -> bool:
    field_count = sum(
        1
        for field in REPORT_FIELDS
        if field in template_text
    )

    section_count = sum(
        1
        for section in REPORT_SECTIONS
        if _normalize_heading(section)
        in _normalize_heading(template_text)
    )

    return field_count >= 3 or section_count >= 3


def _render_report_docx(
    title: str,
    content: str,
    template_text: str,
) -> bytes:
    sections = _template_sections(template_text)

    fields, section_values = _parse_report_content(
        content,
        sections,
    )

    document = DocxDocument()

    heading = document.add_paragraph()
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    run = heading.add_run(
        "업무보고서"
        if "업무보고서" in template_text
        else title
    )
    run.bold = True

    table = document.add_table(
        rows=3,
        cols=4,
    )
    table.style = "Table Grid"

    pairs = [
        (
            "문서명",
            fields.get("문서명", title),
            "작성일",
            fields.get("작성일", ""),
        ),
        (
            "작성부서",
            fields.get("작성부서", ""),
            "작성자",
            fields.get("작성자", ""),
        ),
        (
            "보고대상",
            fields.get("보고대상", ""),
            "문서번호",
            fields.get("문서번호", ""),
        ),
    ]

    for row_index, values in enumerate(pairs):
        label1, value1, label2, value2 = values

        row = table.rows[row_index].cells

        row[0].text = label1
        row[1].text = value1
        row[2].text = label2
        row[3].text = value2

    document.add_paragraph("")

    section_written = False

    for section in sections:
        body = section_values.get(
            section,
            "",
        ).strip()

        if not body:
            continue

        section_written = True

        paragraph = document.add_paragraph()
        run = paragraph.add_run(section)
        run.bold = True

        box = document.add_table(
            rows=1,
            cols=1,
        )
        box.style = "Table Grid"
        box.cell(0, 0).text = body

        document.add_paragraph("")

    if not section_written:
        for line in content.splitlines():
            document.add_paragraph(line)

    buffer = BytesIO()
    document.save(buffer)

    return buffer.getvalue()


def _fill_docx_template(
    title: str,
    content: str,
    template_bytes: bytes,
) -> bytes:
    document = DocxDocument(
        BytesIO(template_bytes)
    )

    template_text = "\n".join(
        [p.text for p in document.paragraphs]
        + [
            cell.text
            for table in document.tables
            for row in table.rows
            for cell in row.cells
        ]
    )

    sections = _template_sections(template_text)

    fields, section_values = _parse_report_content(
        content,
        sections,
    )

    for table in document.tables:
        for row_index, row in enumerate(table.rows):
            for cell_index, cell in enumerate(row.cells):
                label = cell.text.strip()

                if label in fields:
                    if cell_index + 1 < len(row.cells):
                        target = row.cells[cell_index + 1]

                        if not target.text.strip():
                            target.text = fields[label]

                normalized = _normalize_heading(label)

                for section, body in section_values.items():
                    if (
                        normalized
                        == _normalize_heading(section)
                        and body
                    ):
                        if row_index + 1 < len(table.rows):
                            target = table.rows[
                                row_index + 1
                            ].cells[cell_index]

                            if not target.text.strip():
                                target.text = body

    buffer = BytesIO()
    document.save(buffer)

    return buffer.getvalue()


def _generate_template_docx(
    title: str,
    content: str,
    template_bytes: bytes,
    template_filename: str | None,
) -> bytes:
    ext = _file_extension(
        template_filename
    )

    if ext == "docx":
        try:
            return _fill_docx_template(
                title,
                content,
                template_bytes,
            )
        except Exception:
            pass

    template_text = _extract_template_text(
        template_bytes,
        template_filename,
    )

    if (
        template_text
        and _is_report_template(template_text)
    ):
        return _render_report_docx(
            title,
            content,
            template_text,
        )

    return _generate_docx(
        title,
        content,
    )


def _generate_template_pdf(
    title: str,
    content: str,
    template_bytes: bytes,
    template_filename: str | None,
) -> bytes:
    template_text = _extract_template_text(
        template_bytes,
        template_filename,
    )

    if not (
        template_text
        and _is_report_template(template_text)
    ):
        return _generate_pdf(
            title,
            content,
        )

    sections = _template_sections(
        template_text
    )

    fields, section_values = _parse_report_content(
        content,
        sections,
    )

    lines = [
        "업무보고서",
        "",
        f"문서명: {fields.get('문서명', title)}",
        f"작성일: {fields.get('작성일', '')}",
        f"작성부서: {fields.get('작성부서', '')}",
        f"작성자: {fields.get('작성자', '')}",
        f"보고대상: {fields.get('보고대상', '')}",
        f"문서번호: {fields.get('문서번호', '')}",
        "",
    ]

    for section in sections:
        body = section_values.get(
            section,
            "",
        ).strip()

        if body:
            lines.extend(
                [
                    section,
                    body,
                    "",
                ]
            )

    return _generate_pdf(
        title,
        "\n".join(lines),
    )


def _generate_xlsx(
    title: str,
    content: str,
    template_bytes: bytes | None = None,
    template_filename: str | None = None,
) -> bytes:
    template_text = ""

    if template_bytes:
        template_text = _extract_template_text(
            template_bytes,
            template_filename,
        )

    sections = _template_sections(
        template_text
    )

    fields, section_values = _parse_report_content(
        content,
        sections,
    )

    if (
        template_bytes
        and _file_extension(template_filename)
        == "xlsx"
    ):
        workbook = load_workbook(
            BytesIO(template_bytes)
        )

        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    value = (
                        str(cell.value).strip()
                        if cell.value is not None
                        else ""
                    )

                    if value in fields:
                        target = sheet.cell(
                            row=cell.row,
                            column=cell.column + 1,
                        )

                        if target.value in (None, ""):
                            target.value = fields[value]

        buffer = BytesIO()
        workbook.save(buffer)

        return buffer.getvalue()

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "보고서"

    sheet["A1"] = title

    row_index = 3

    for field in REPORT_FIELDS:
        if field in fields:
            sheet.cell(
                row=row_index,
                column=1,
                value=field,
            )
            sheet.cell(
                row=row_index,
                column=2,
                value=fields[field],
            )
            row_index += 1

    row_index += 1

    for section in sections:
        body = section_values.get(
            section,
            "",
        ).strip()

        if not body:
            continue

        sheet.cell(
            row=row_index,
            column=1,
            value=section,
        )

        sheet.cell(
            row=row_index + 1,
            column=1,
            value=body,
        )

        row_index += 3

    buffer = BytesIO()
    workbook.save(buffer)

    return buffer.getvalue()


def _generate_pptx(title: str, content: str) -> bytes:
    """최소 기능 버전: 제목 슬라이드 1장 + 본문을 문단 단위로 나눠 슬라이드에 배치.
    (디자인 템플릿, 이미지, 차트 등은 범위 밖 - 추후 고도화 필요)"""
    prs = Presentation()

    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    # 본문 - 빈 줄 기준으로 문단을 나눠 슬라이드마다 하나씩
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    content_layout = prs.slide_layouts[1]  # 제목 + 본문 레이아웃

    for paragraph in paragraphs:
        slide = prs.slides.add_slide(content_layout)
        lines = paragraph.split("\n")
        slide.shapes.title.text = lines[0][:60]  # 첫 줄을 슬라이드 제목으로

        if len(lines) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = lines[1]
            for line in lines[2:]:
                p = tf.add_paragraph()
                p.text = line

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _generate_markdown(
    title: str,
    content: str,
) -> bytes:
    return (
        f"# {title}\n\n"
        f"{content.strip()}\n"
    ).encode("utf-8")


def _generate_text(
    content: str,
) -> bytes:
    return content.encode("utf-8")

def generate_document(
    title: str,
    content: str,
    output_format: str,
    template_bytes: bytes | None = None,
    template_filename: str | None = None,
) -> tuple[bytes, str, str]:
    fmt = output_format.lower().strip()
    safe_title = _safe_filename(title)

    if fmt == "docx":
        data = (
            _generate_template_docx(
                title,
                content,
                template_bytes,
                template_filename,
            )
            if template_bytes
            else _generate_docx(
                title,
                content,
            )
        )

        return (
            data,
            f"{safe_title}.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    if fmt == "pdf":
        data = (
            _generate_template_pdf(
                title,
                content,
                template_bytes,
                template_filename,
            )
            if template_bytes
            else _generate_pdf(
                title,
                content,
            )
        )

        return (
            data,
            f"{safe_title}.pdf",
            "application/pdf",
        )

    if fmt == "xlsx":
        return (
            _generate_xlsx(
                title,
                content,
                template_bytes,
                template_filename,
            ),
            f"{safe_title}.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    if fmt == "pptx":
        return (
            _generate_pptx(title, content),
            f"{safe_title}.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    if fmt == "md":
        return (
            _generate_markdown(
                title,
                content,
            ),
            f"{safe_title}.md",
            "text/markdown; charset=utf-8",
        )

    if fmt == "txt":
        return (
            _generate_text(content),
            f"{safe_title}.txt",
            "text/plain; charset=utf-8",
        )

    raise ValueError(
        "지원 형식은 docx, pdf, xlsx, pptx, md, txt입니다."
    )
